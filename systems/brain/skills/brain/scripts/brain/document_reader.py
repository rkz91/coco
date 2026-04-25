"""
Brain Document Reader — Extracts plain text from office documents.

Supports: .docx, .pdf, .xlsx, .pptx
Returns plain text suitable for entity extraction and brain ingestion.
All readers are best-effort — return empty string on failure, never raise.
"""

import logging
from pathlib import Path

log = logging.getLogger(__name__)


def read_docx(path: str | Path) -> str:
    """Extract text from a .docx file. Returns plain text or empty string."""
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        # Also read tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        log.debug("read_docx failed for %s: %s", path, e)
        return ""


def read_pdf(path: str | Path) -> str:
    """Extract text from a .pdf file. Returns plain text or empty string."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages[:50]:  # Cap at 50 pages
            text = page.extract_text()
            if text and text.strip():
                parts.append(text.strip())
        return "\n\n".join(parts)
    except Exception as e:
        log.debug("read_pdf failed for %s: %s", path, e)
        return ""


def read_xlsx(path: str | Path) -> str:
    """Extract text from an .xlsx file. Returns plain text or empty string."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames[:10]:  # Cap at 10 sheets
            ws = wb[sheet]
            parts.append(f"=== Sheet: {sheet} ===")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
                    row_count += 1
                    if row_count > 200:  # Cap rows per sheet
                        parts.append("[... truncated ...]")
                        break
        wb.close()
        return "\n".join(parts)
    except Exception as e:
        log.debug("read_xlsx failed for %s: %s", path, e)
        return ""


def read_pptx(path: str | Path) -> str:
    """Extract text from a .pptx file. Returns plain text or empty string."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                parts.append(f"--- Slide {i+1} ---")
                parts.extend(slide_texts)
        return "\n".join(parts)
    except Exception as e:
        log.debug("read_pptx failed for %s: %s", path, e)
        return ""


def read_document(path: str | Path) -> str:
    """Auto-detect file type and extract text. Returns plain text or empty string."""
    path = Path(path)
    suffix = path.suffix.lower()
    readers = {
        ".docx": read_docx,
        ".pdf": read_pdf,
        ".xlsx": read_xlsx,
        ".xls": read_xlsx,  # openpyxl handles some .xls
        ".pptx": read_pptx,
    }
    reader = readers.get(suffix)
    if reader:
        return reader(path)
    # For text-based files, read directly
    if suffix in (".md", ".html", ".txt", ".csv", ".json", ".xml", ".yml", ".yaml"):
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return ""
    return ""


def is_readable(path: str | Path) -> bool:
    """Check if a file can be read by the document reader."""
    suffix = Path(path).suffix.lower()
    return suffix in {".docx", ".pdf", ".xlsx", ".xls", ".pptx", ".md", ".html", ".txt", ".csv", ".json", ".xml", ".yml", ".yaml"}
